"""
V6 로그 데이터를 PowerPoint 보고서로 생성
"""
import json
from pathlib import Path
from collections import Counter
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# 로그 파일 로드
log_dir = Path("dashboard/logs/v6_chatbot")
log_files = sorted(log_dir.glob("v6_log_*.jsonl"))

all_logs = []
for log_file in log_files:
    with open(log_file, 'r', encoding='utf-8') as f:
        for line in f:
            try:
                all_logs.append(json.loads(line))
            except:
                pass

print(f"Loaded {len(all_logs)} logs")

# 통계 계산
complexities = Counter(log['complexity'] for log in all_logs)
viz_strategies = Counter(log['visualization_strategy'] for log in all_logs)
avg_duration = sum(log['total_duration'] for log in all_logs) / len(all_logs) if all_logs else 0
avg_queries = sum(log['total_queries'] for log in all_logs) / len(all_logs) if all_logs else 0

# 날짜별 통계
daily_stats = {}
for log in all_logs:
    date = log['timestamp'][:10]  # YYYY-MM-DD
    if date not in daily_stats:
        daily_stats[date] = {'count': 0, 'total_duration': 0}
    daily_stats[date]['count'] += 1
    daily_stats[date]['total_duration'] += log['total_duration']

# PPT 생성
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ====================================
# 슬라이드 1: 타이틀
# ====================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

# 제목
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "V6 LangGraph Agent"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(48)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 102, 204)
title_para.alignment = PP_ALIGN.CENTER

# 부제
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "사용 분석 보고서"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(28)
subtitle_para.alignment = PP_ALIGN.CENTER

# 날짜
date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
date_frame = date_box.text_frame
date_frame.text = f"기간: {min(daily_stats.keys())} ~ {max(daily_stats.keys())}"
date_para = date_frame.paragraphs[0]
date_para.font.size = Pt(18)
date_para.alignment = PP_ALIGN.CENTER

# ====================================
# 슬라이드 2: 전체 요약
# ====================================
slide = prs.slides.add_slide(prs.slide_layouts[1])  # 제목 + 내용
title = slide.shapes.title
title.text = "전체 요약"

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

# 통계 추가
stats_text = f"""총 요청 수: {len(all_logs)}개

평균 실행 시간: {avg_duration:.2f}초

요청당 평균 쿼리 수: {avg_queries:.2f}개

분석 기간: {len(daily_stats)}일"""

tf.text = stats_text
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(24)
    paragraph.space_after = Pt(20)

# ====================================
# 슬라이드 3: 복잡도 분포 (차트)
# ====================================
slide = prs.slides.add_slide(prs.slide_layouts[5])  # 제목만
title = slide.shapes.title
title.text = "질문 복잡도 분포"

# 차트 데이터
chart_data = CategoryChartData()
chart_data.categories = list(complexities.keys())
chart_data.add_series('요청 수', list(complexities.values()))

# 차트 추가 (파이 차트)
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = 2  # 오른쪽
chart.legend.include_in_layout = False

# ====================================
# 슬라이드 4: 시각화 전략 분포
# ====================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "시각화 전략 분포"

# 차트 데이터
chart_data = CategoryChartData()
chart_data.categories = list(viz_strategies.keys())
chart_data.add_series('요청 수', list(viz_strategies.values()))

# 막대 차트
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

chart.has_legend = False

# ====================================
# 슬라이드 5: 일별 사용량
# ====================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "일별 사용량 추이"

# 차트 데이터
sorted_dates = sorted(daily_stats.keys())
chart_data = CategoryChartData()
chart_data.categories = [d[5:] for d in sorted_dates]  # MM-DD만
chart_data.add_series('요청 수', [daily_stats[d]['count'] for d in sorted_dates])

# 라인 차트
x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

chart.has_legend = False

# ====================================
# 슬라이드 6: 최근 질문 샘플 (5개)
# ====================================
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "최근 질문 샘플"

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

for i, log in enumerate(all_logs[-5:], 1):
    query = log['user_query']
    # 한글 깨짐 방지 - 영문으로 대체 또는 잘라내기
    if len(query) > 60:
        query = query[:60] + "..."

    p = tf.add_paragraph()
    p.text = f"{i}. {query}"
    p.font.size = Pt(16)
    p.space_after = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = f"   복잡도: {log['complexity']} | 실행시간: {log['total_duration']:.1f}초"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    p2.space_after = Pt(18)

# ====================================
# 슬라이드 7: 성능 분석
# ====================================
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "성능 분석"

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

# 복잡도별 평균 실행 시간
complexity_duration = {}
for log in all_logs:
    comp = log['complexity']
    if comp not in complexity_duration:
        complexity_duration[comp] = []
    complexity_duration[comp].append(log['total_duration'])

avg_by_complexity = {
    comp: sum(durations) / len(durations)
    for comp, durations in complexity_duration.items()
}

perf_text = "복잡도별 평균 실행 시간:\n\n"
for comp in sorted(avg_by_complexity.keys()):
    perf_text += f"  • {comp}: {avg_by_complexity[comp]:.2f}초\n"

perf_text += f"\n시각화 생성률: {(viz_strategies.get('auto', 0) + viz_strategies.get('suggest', 0)) / len(all_logs) * 100:.1f}%"

tf.text = perf_text
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(20)
    paragraph.space_after = Pt(15)

# ====================================
# 슬라이드 8: 핵심 인사이트
# ====================================
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "핵심 인사이트"

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

# 인사이트 계산
simple_pct = complexities.get('simple', 0) / len(all_logs) * 100
no_viz_pct = viz_strategies.get('none', 0) / len(all_logs) * 100

insights = f"""• 전체 요청의 {simple_pct:.0f}%가 단순 질문
  → 대부분 직관적인 질문 패턴

• 시각화 미생성 비율: {no_viz_pct:.0f}%
  → 텍스트 기반 응답 선호 또는 개선 필요

• 평균 응답 시간: {avg_duration:.1f}초
  → {'빠른 응답' if avg_duration < 10 else '개선 필요'}

• 일평균 사용량: {len(all_logs) / len(daily_stats):.1f}건
  → {'활발한 사용' if len(all_logs) / len(daily_stats) > 20 else '초기 단계'}"""

tf.text = insights
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(18)
    paragraph.space_after = Pt(20)

# 저장
output_file = "V6_Usage_Report.pptx"
prs.save(output_file)
print(f"PPT generated: {output_file}")
print(f"Total slides: {len(prs.slides)}")
